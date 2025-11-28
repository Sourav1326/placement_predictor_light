import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_placement_predictor_presentation():
    """Create a PowerPoint presentation for the Placement Predictor project"""
    
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Define styles
    def set_title_style(title):
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def set_subtitle_style(subtitle):
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)
        subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def set_content_style(content):
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Industry-Ready Placement Prediction System"
    subtitle.text = "A Comprehensive Career Guidance Platform\n\nAdvanced Machine Learning & Deep Learning\nSkill Verification & Assessment\nPersonalized Recommendations\nATS-Optimized Resume Analysis"
    
    # Slide 2: Project Overview
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Project Overview"
    content.text = ("What is the Placement Predictor?\n\n"
                   "A complete career guidance platform designed to help students improve their job placement prospects through:\n\n"
                   "• Data-Driven Insights: Advanced analytics for accurate placement predictions\n"
                   "• Skill Verification: Real-time assessment of technical and soft skills\n"
                   "• AI-Powered Recommendations: Personalized career development paths\n"
                   "• Comprehensive Assessment: Multi-dimensional evaluation of student capabilities\n\n"
                   "Target Users:\n"
                   "• Students seeking accurate placement predictions\n"
                   "• Educational institutions aiming to increase placement rates\n"
                   "• Placement officers needing analytics and tracking tools")
    
    # Slide 3: Core Problems Addressed
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Core Problems Addressed"
    content.text = ("Why This Project Matters\n\n"
                   "Traditional placement systems face several challenges:\n\n"
                   "❌ Low Accuracy: Traditional systems often provide inaccurate predictions\n"
                   "❌ Unverified Skills: Employers can't trust self-reported skills\n"
                   "❌ Poor ATS Compatibility: Resumes fail to pass automated screening\n"
                   "❌ Generic Guidance: One-size-fits-all career advice\n"
                   "❌ Skill Mismatch: Students lack industry-required competencies\n\n"
                   "✅ Our Solution: Trust-weighted ML predictions using verified skills\n"
                   "✅ Our Solution: Real-time skill verification with badge system\n"
                   "✅ Our Solution: ATS-optimized resume analysis and autofill\n"
                   "✅ Our Solution: Proactive AI guidance based on user behavior\n"
                   "✅ Our Solution: Multi-model ensemble approach for higher accuracy")
    
    # Slide 4: Key Features
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Key Features"
    content.text = ("Comprehensive Feature Set\n\n"
                   "🔍 Skill Verification System\n"
                   "• Live coding challenges with real-time evaluation\n"
                   "• SQL sandbox for database skills testing\n"
                   "• Framework code review with 4-level badge system\n"
                   "• Light proctoring for integrity assurance\n\n"
                   "📄 ATS Resume Analyzer\n"
                   "• Smart autofill based on profile data\n"
                   "• Compatibility scoring with detailed feedback\n"
                   "• Interactive fixes for optimization\n"
                   "• Industry-specific template recommendations\n\n"
                   "🤖 AI Career Chatbot\n"
                   "• Context-aware conversations with personalized suggestions\n"
                   "• Proactive career guidance based on user behavior\n"
                   "• Instant action buttons for immediate engagement\n"
                   "• 24/7 availability for student support")
    
    # Slide 5: Technical Architecture
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technical Architecture"
    content.text = ("System Architecture Overview\n\n"
                   "Frontend (Bootstrap) ↔ Flask Web Server ↔ Business Logic ↔ ML Models & NLP ↔ SQLite Database\n\n"
                   "Design Patterns Used:\n"
                   "• MVC-like Pattern: Templates, Flask routes, backend modules\n"
                   "• Singleton Pattern: Database connection management\n"
                   "• Factory Pattern: Model training and assessment engines\n"
                   "• Observer Pattern: Chatbot notifications\n\n"
                   "Component Interaction:\n"
                   "• Frontend ↔ Flask server ↔ Business logic modules ↔ ML models & NLP processors ↔ SQLite database\n"
                   "• External integrations via REST APIs and webhook support\n"
                   "• Sandboxed environments for secure code execution")
    
    # Slide 6: Technology Stack
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technology Stack"
    content.text = ("Modern Technology Ecosystem\n\n"
                   "🖥️ Frontend\n"
                   "• Bootstrap UI with responsive design\n"
                   "• Jinja2 templating for dynamic content\n"
                   "• Interactive dashboards and visualizations\n\n"
                   "⚙️ Backend\n"
                   "• Flask 3.0.0 as the web framework\n"
                   "• Werkzeug 3.0.1 for WSGI utilities\n"
                   "• SQLite for lightweight database management\n\n"
                   "🧠 Machine Learning\n"
                   "• scikit-learn 1.3.2 for traditional ML models\n"
                   "• XGBoost for gradient boosting algorithms\n"
                   "• TensorFlow 2.15.0 & Keras for deep learning\n"
                   "• PyTorch 2.1.1 for neural networks")
    
    # Slide 7: Machine Learning Models
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Machine Learning Models"
    content.text = ("Advanced Prediction Algorithms\n\n"
                   "📈 Ensemble Approach\n"
                   "Multiple models working together for higher accuracy:\n\n"
                   "1. Logistic Regression: Baseline model with hyperparameter tuning\n"
                   "2. Random Forest: Tree-based ensemble with feature importance\n"
                   "3. XGBoost: Gradient boosting for complex pattern recognition\n"
                   "4. Deep Neural Network: Multi-layer perceptron for non-linear relationships\n\n"
                   "🎯 Model Performance\n"
                   "• Cross-validation for robust evaluation\n"
                   "• AUC-ROC as primary evaluation metric\n"
                   "• Feature importance analysis for interpretability\n"
                   "• Continuous learning from new data")
    
    # Slide 8: Deep Learning Implementation
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Deep Learning Implementation"
    content.text = ("Neural Network Architecture\n\n"
                   "🧠 Advanced Deep Learning Model\n"
                   "• Multi-layer neural network with dropout regularization\n"
                   "• Batch normalization for stable training\n"
                   "• Adam optimizer with learning rate scheduling\n"
                   "• Early stopping to prevent overfitting\n\n"
                   "📊 Architecture Details\n"
                   "• Input layer: 34 features (academic, technical, soft skills)\n"
                   "• Hidden layers: 256 → 128 → 64 → 32 neurons with ReLU activation\n"
                   "• Output layer: Single neuron with sigmoid activation\n"
                   "• Regularization: L2 regularization and dropout layers\n\n"
                   "🎯 Training Results\n"
                   "• Validation AUC: 1.0000 (on test data)\n"
                   "• Validation Accuracy: 0.9500 (on test data)")
    
    # Slide 9: Assessment & Verification System
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Assessment & Verification System"
    content.text = ("Comprehensive Evaluation Framework\n\n"
                   "🧪 Multi-Dimensional Assessments\n"
                   "1. Comprehensive Aptitude Tests: Cognitive and analytical skills\n"
                   "2. Technical Skill Quizzes: Programming language proficiency\n"
                   "3. Communication Analysis: Written and verbal communication\n"
                   "4. Situational Judgment Tests: Professional decision-making\n"
                   "5. Mock Interviews: Realistic interview simulation\n\n"
                   "🔐 Trust but Verify Approach\n"
                   "• 4-Level Badge System: Basic → Intermediate → Advanced → Verified\n"
                   "• Live Coding Challenges: Real-time code evaluation\n"
                   "• SQL Sandbox: Database query assessment\n"
                   "• Framework Code Review: Best practices verification\n"
                   "• Light Proctoring: Integrity assurance without invasion")
    
    # Slide 10: Career Guidance Features
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Career Guidance Features"
    content.text = ("Personalized Career Development\n\n"
                   "🤖 AI Career Chatbot\n"
                   "• Context-aware conversations with personalized suggestions\n"
                   "• Proactive career guidance based on user behavior\n"
                   "• Instant action buttons for immediate engagement\n"
                   "• 24/7 availability for student support\n\n"
                   "📚 Course Recommendation Engine\n"
                   "• Personalized learning path recommendations\n"
                   "• Free and paid course alternatives\n"
                   "• Company-specific skill development paths\n"
                   "• Progress tracking and completion certificates\n\n"
                   "🔍 Smart Search Engine\n"
                   "• Intelligent job search with filters\n"
                   "• Skill-based matching algorithms\n"
                   "• Salary range predictions\n"
                   "• Location-based opportunities")
    
    # Slide 11: Portable Environment
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Portable Environment"
    content.text = ("Offline-Ready Deployment\n\n"
                   "📦 Self-Contained Environment\n"
                   "• All dependencies pre-installed in portable environment\n"
                   "• No internet required after initial setup\n"
                   "• No repeated downloads or installations\n"
                   "• Consistent environment across different machines\n\n"
                   "🚀 Quick Start Options\n"
                   "1. Direct Run: Double-click run_direct.bat\n"
                   "2. Development Mode: Activate environment and run manually\n"
                   "3. Production Deployment: Using Gunicorn or Waitress\n\n"
                   "📁 Directory Structure\n"
                   "placement predictor/\n"
                   "├── portable_env/           # Python environment\n"
                   "├── requirements_cache/     # Cached packages\n"
                   "├── src/                    # Source code\n"
                   "├── data/                   # Data files\n"
                   "├── templates/              # HTML templates\n"
                   "└── models/                 # Trained models")
    
    # Slide 12: Database Schema
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Database Schema"
    content.text = ("Comprehensive Data Management\n\n"
                   "🗃️ Key Tables\n"
                   "1. Users: Student and admin account management\n"
                   "2. Student Profiles: Detailed academic and skill information\n"
                   "3. Assessment Results: Skill evaluation tracking\n"
                   "4. Placement Predictions: Prediction history and analytics\n"
                   "5. Course Progress: Learning path tracking\n"
                   "6. User Sessions: Authentication and session management\n\n"
                   "🔗 Relationships\n"
                   "• One-to-one: Users ↔ Student Profiles\n"
                   "• One-to-many: Users ↔ Assessments, Predictions, Course Progress\n"
                   "• Foreign key constraints for data integrity\n"
                   "• Cascading deletes for clean data management")
    
    # Slide 13: Implementation & Results
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Implementation & Results"
    content.text = ("Project Success Metrics\n\n"
                   "🎯 Key Achievements\n"
                   "• Portable Environment: Fully self-contained with offline capability\n"
                   "• Model Performance: High accuracy with ensemble approach\n"
                   "• Skill Verification: Real-time assessment with badge system\n"
                   "• User Experience: Intuitive interface with comprehensive features\n\n"
                   "📊 Performance Benchmarks\n"
                   "• Prediction response time: < 200ms\n"
                   "• Database queries: < 50ms\n"
                   "• Code execution: < 5 seconds\n"
                   "• Supports 500+ concurrent users\n\n"
                   "✅ Quality Assurance\n"
                   "• Comprehensive testing of all modules\n"
                   "• Cross-validation for model robustness\n"
                   "• Error handling and graceful degradation\n"
                   "• Security measures for data protection")
    
    # Slide 14: Future Enhancements
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Future Enhancements"
    content.text = ("Roadmap for Continuous Improvement\n\n"
                   "🚀 Planned Features\n"
                   "1. Mobile Application: Native mobile experience for on-the-go access\n"
                   "2. Industry Partnerships: Direct integration with company recruitment systems\n"
                   "3. Advanced Analytics: Predictive analytics for institutional planning\n"
                   "4. Multi-Language Support: Global accessibility and localization\n"
                   "5. Blockchain Verification: Immutable skill verification records\n\n"
                   "🛠️ Technical Improvements\n"
                   "• Enhanced deep learning architectures\n"
                   "• Real-time model updating with streaming data\n"
                   "• Improved natural language processing capabilities\n"
                   "• Advanced visualization and reporting tools\n"
                   "• Scalability enhancements for enterprise deployment")
    
    # Slide 15: Conclusion
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Conclusion"
    content.text = ("Transforming Career Development\n\n"
                   "🌟 Key Benefits\n"
                   "• For Students: Accurate predictions, skill verification, personalized guidance\n"
                   "• For Institutions: Higher placement rates, analytics, student tracking\n"
                   "• For Employers: Verified skills, better candidate matching\n\n"
                   "🎯 Impact\n"
                   "• Increased placement success rates\n"
                   "• Reduced time-to-hire for employers\n"
                   "• Enhanced student confidence and preparedness\n"
                   "• Data-driven career development decisions\n\n"
                   "🚀 Ready for Deployment\n"
                   "• Fully functional portable environment\n"
                   "• Comprehensive documentation\n"
                   "• Easy setup and maintenance\n"
                   "• Scalable architecture for growth\n\n"
                   "Thank You!\n"
                   "Questions & Discussion")
    
    # Slide 16: Demo & Q&A
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Demo & Q&A"
    content.text = ("Live Demonstration\n\n"
                   "🖥️ Quick Start Guide\n"
                   "1. Double-click run_direct.bat\n"
                   "2. Open browser to http://localhost:5000\n"
                   "3. Login with admin credentials:\n"
                   "   • Email: admin@placement.system\n"
                   "   • Password: admin123\n\n"
                   "📋 Key Demo Features\n"
                   "• Student dashboard with placement predictions\n"
                   "• Skill assessment and verification\n"
                   "• AI chatbot interaction\n"
                   "• Resume analysis and optimization\n"
                   "• Course recommendations\n\n"
                   "❓ Questions & Answers\n"
                   "Open floor for questions and discussion")
    
    # Save presentation
    prs.save("Placement_Predictor_Presentation.pptx")
    print("Presentation created successfully as 'Placement_Predictor_Presentation.pptx'")

if __name__ == "__main__":
    create_placement_predictor_presentation()